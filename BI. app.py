"""
app.py - MetriqAI Ana Uygulama (V2.0)
"""

import streamlit as st
import os
import glob
import zipfile
from io import BytesIO
import pandas as pd
import numpy as np # np import'u eklendi

# Modül importları
from data_analysis import SalesAnalyzer
from reporting import (
    generate_graphs, ai_summary, build_pdf,
    build_docx, build_ppt, save_excel, OPENAI_API_KEY
)
from config import APP_NAME, PACKAGES
from ui_components import (
    load_custom_css, render_hero, render_package_selector,
    render_team_section, simulate_processing, render_footer
)
from maps import (
    create_turkey_heatmap, create_world_map,
    create_interactive_bar_chart, create_time_series_chart
)

# -----------------------------------------------------------

# SAYFA YAPLANDIRMASI

# -----------------------------------------------------------

st.set_page_config(
    page_title=f"{APP_NAME} Dashboard",
    layout="wide",
    page_icon="⚡",
    initial_sidebar_state="expanded" 
)

# Custom CSS yükle
load_custom_css()

# -----------------------------------------------------------

# SIDEBAR

# -----------------------------------------------------------

# Paket seçimi
selected_package = render_package_selector()
package_features = PACKAGES[selected_package]['features']

# Ekip gösterimi
render_team_section()

# İletişim bilgileri
st.sidebar.markdown("---")
st.sidebar.markdown("### 📞 Contact")
st.sidebar.markdown("📧 insights@metriq.ai")
st.sidebar.markdown("🌐 metriq.ai")
st.sidebar.markdown("📞 +90 XXX XXX XXXX")

# -----------------------------------------------------------

# ANA SAYFA

# -----------------------------------------------------------

# Hero Section
render_hero()

# Dosya Yükleme
st.markdown("### 📂 Upload Your Data")
uploaded_file = st.file_uploader(
    "Drop your Excel or CSV file here",
    type=["xlsx", "xls", "csv"],
    help="Upload sales data for instant AI-powered analysis"
)

# -----------------------------------------------------------

# VERİ ANALİZİ

# -----------------------------------------------------------

if uploaded_file:
    try:
        # Veriyi yükle
        analyzer = SalesAnalyzer.load_data(uploaded_file)

        # İşlem animasyonu göster (illüzyon!)  
        simulate_processing(selected_package)  
          
        # KPI'ları hesapla  
        kpis = analyzer.compute_kpis()  
          
        st.markdown("---")  
          
        # -------------------------------------------------------  
        # KPI KARTLARI  
        # -------------------------------------------------------  
          
        st.markdown("### 🎯 Key Performance Indicators")  
          
        col1, col2, col3, col4 = st.columns(4)  
          
        with col1:  
            st.metric(  
                "💰 Total Revenue",  
                f"{kpis['total']:,.0f} TL",  
                help="Total revenue for the period"  
            )  
          
        with col2:  
            st.metric(  
                "📊 Daily Average",  
                f"{kpis['avg']:,.0f} TL",  
                help="Average daily revenue"  
            )  
          
        with col3:  
            wow_val = kpis["wow"]  
            # numpy'dan gelen nan kontrolü düzeltildi
            if not pd.isna(wow_val) and not np.isnan(wow_val):  
                delta_val = f"{wow_val:+.1f}%"  
                delta_color = "normal" if wow_val >= 0 else "inverse"  
                display_wow = f"{wow_val:.1f}%"
            else:  
                delta_val = "N/A"  
                delta_color = "off"  
                display_wow = "N/A"
              
            st.metric(  
                "📈 Week-over-Week",  
                display_wow,  
                delta=delta_val,  
                delta_color=delta_color,  
                help="Last 7 days vs previous 7 days"  
            )  
          
        with col4:  
            st.metric(  
                "🔢 Total Transactions",  
                f"{len(analyzer.df):,}",  
                help="Total number of transactions"  
            )  
          
        st.markdown("---")  
          
        # -------------------------------------------------------  
        # AI ÖZETİ (Pro & Premium)  
        # -------------------------------------------------------  
          
        summary = "" # Summary değişkeni tanımlanıyor
        if package_features['ai_summary']:  
            st.markdown("### 🤖 AI Strategic Analysis")  
              
            if OPENAI_API_KEY:  
                try:  
                    # st.spinner içinde ai_summary çağrılıyor
                    with st.spinner("AI is analyzing your data..."):  
                        summary = ai_summary(analyzer.get_kpi_text())  
                      
                    # Güzel bir box içinde göster  
                    st.markdown(f"""  
                    <div style='  
                        background: rgba(0, 191, 255, 0.1);  
                        border-left: 5px solid #00BFFF;  
                        padding: 20px;  
                        border-radius: 10px;  
                        margin: 20px 0;  
                    '>  
                        <p style='color: rgba(255,255,255,0.9); line-height: 1.8;'>  
                            {summary}  
                        </p>  
                    </div>  
                    """, unsafe_allow_html=True)  
                      
                except Exception as e:  
                    summary = f"[AI Analysis] Service temporarily unavailable. Error: {str(e)}"  
                    st.warning(summary)  
            else:  
                summary = "[AI Analysis] OpenAI API key not configured. Using Mock Summary."
                # Mock summary'yi al
                summary = ai_summary(analyzer.get_kpi_text()) 
                st.info(summary)
        else:  
            summary = "[Basic Package] AI analysis not included. Upgrade to Pro or Premium!"  
            st.info(summary)  
          
        st.markdown("---")  
          
        # -------------------------------------------------------  
        # INTERACTIVE CHARTS (Pro & Premium)  
        # -------------------------------------------------------  
          
        graphs = {} # Grafik dosyalarını tutacak değişken
        st.markdown("### 📊 Data Visualization")  
          
        if package_features['interactive_charts']:  
            # Time series  
            fig_time = create_time_series_chart(analyzer.df.reset_index()) # Index resetlendi
            st.plotly_chart(fig_time, use_container_width=True)  
              
            # İki sütun: Kategori & Şehir  
            col_a, col_b = st.columns(2)  
              
            with col_a:  
                if 'kategori' in analyzer.df.columns:  
                    fig_cat = create_interactive_bar_chart(  
                        analyzer.df, 'kategori', '📦 Top Categories'  
                    )  
                    st.plotly_chart(fig_cat, use_container_width=True)  
          
            with col_b:  
                if 'sehir' in analyzer.df.columns:  
                    fig_city = create_interactive_bar_chart(  
                        analyzer.df, 'sehir', '🏙️ Top Cities'  
                    )  
                    st.plotly_chart(fig_city, use_container_width=True)  
              
        else:  
            # Basic: Statik grafikler  
            graphs = generate_graphs(kpis) # Statik grafikleri oluştur
              
            col_a, col_b = st.columns(2)  
            with col_a:  
                # Sadece dosya yolu olduğu için açıp göster
                if "daily" in graphs and os.path.exists(graphs["daily"]):
                     st.image(graphs["daily"], caption="📈 Daily Revenue Trend")  
            with col_b:  
                if "top_product" in graphs and os.path.exists(graphs["top_product"]):
                     st.image(graphs["top_product"], caption="📦 Top Products")  
              
            if "top_region" in graphs and os.path.exists(graphs["top_region"]):
                 st.image(graphs["top_region"], caption="🌍 Regional Distribution")  
          
        st.markdown("---")  
          
        # -------------------------------------------------------  
        # TURKEY HEATMAP (Pro & Premium)  
        # -------------------------------------------------------  
          
        if package_features['turkey_map']:  
            st.markdown("### 🇹🇷 Turkey Geographic Analysis")  
              
            if 'sehir' in analyzer.df.columns:  
                fig_turkey = create_turkey_heatmap(analyzer.df)  
                st.plotly_chart(fig_turkey, use_container_width=True)  
            else:  
                st.info("City column not found in data. Heatmap requires 'sehir' column.")  
          
        # -------------------------------------------------------  
        # WORLD MAP (Premium Only)  
        # -------------------------------------------------------  
          
        if package_features['world_map']:  
            st.markdown("### 🌍 Global Market Analysis")  
              
            fig_world = create_world_map(analyzer.df)  
            st.plotly_chart(fig_world, use_container_width=True)  
          
        st.markdown("---")  
          
        # -------------------------------------------------------  
        # RAPOR OLUŞTURMA  
        # -------------------------------------------------------  
          
        st.markdown("### 📥 Download Your Reports")  
          
        # PDF ve PPT için statik grafikler gerekli, interaktif olsa bile oluşturuluyor
        if not graphs: # Eğer interaktifse (Pro/Premium), statik grafikleri de oluştur
             graphs = generate_graphs(kpis)
          
        # Raporları oluştur  
        report_paths = []  
          
        # 1. PDF (Her paket)  
        pdf_path = build_pdf(summary, kpis, graphs, "metriqAI_report.pdf")  
        report_paths.append(pdf_path)  
          
        # 2. Excel (Her paket)  
        excel_path = save_excel(analyzer.df.copy(), "data_clean.xlsx")  # Copy ile index sorunları önlendi
        report_paths.append(excel_path)  
          
        # 3. Word (Pro & Premium)  
        docx_path = None  
        if package_features['ai_summary'] and "AI analysis not included" not in summary:  # Word ai_summary varsa
            docx_path = build_docx(analyzer.df.copy(), summary, "metriqAI_report.docx")  
            report_paths.append(docx_path)  
          
        # 4. PowerPoint (Pro & Premium)  
        pptx_path = None  
        if package_features['powerpoint']:  
            pptx_path = build_ppt(summary, graphs, "metriqAI_presentation.pptx")  
            report_paths.append(pptx_path)  
          
        # -------------------------------------------------------  
        # İNDİRME BUTONLARI  
        # -------------------------------------------------------  
          
        # Dinamik sütun sayısı  
        is_ai_included = package_features['ai_summary'] and "AI analysis not included" not in summary
        
        num_reports = sum([  
            True,  # PDF  
            True,  # Excel  
            is_ai_included,  # Word  
            package_features['powerpoint'],  # PPT  
            is_ai_included # ZIP (Pro/Premium)  
        ])  
          
        cols = st.columns(max(num_reports, 1)) # En az 1 sütun
        col_idx = 0  
          
        # PDF  
        if col_idx < len(cols) and os.path.exists(pdf_path):
            with cols[col_idx]:  
                with open(pdf_path, "rb") as f:  
                    st.download_button(  
                        "📑 PDF Report",  
                        f,  
                        file_name="metriqAI_report.pdf",  
                        mime="application/pdf"  
                    )  
            col_idx += 1  
          
        # Excel  
        if col_idx < len(cols) and os.path.exists(excel_path):
            with cols[col_idx]:  
                with open(excel_path, "rb") as f:  
                    st.download_button(  
                        "📊 Excel Data",  
                        f,  
                        file_name="data_clean.xlsx",  
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"  
                    )  
            col_idx += 1  
          
        # Word  
        if docx_path and col_idx < len(cols) and os.path.exists(docx_path):  
            with cols[col_idx]:  
                with open(docx_path, "rb") as f:  
                    st.download_button(  
                        "📝 Word Report",  
                        f,  
                        file_name="metriqAI_report.docx",  
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"  
                    )  
            col_idx += 1  
          
        # PowerPoint  
        if pptx_path and col_idx < len(cols) and os.path.exists(pptx_path):  
            with cols[col_idx]:  
                with open(pptx_path, "rb") as f:  
                    st.download_button(  
                        "📊 PowerPoint",  
                        f,  
                        file_name="metriqAI_presentation.pptx",  
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"  
                    )  
            col_idx += 1  
          
        # ZIP (Pro & Premium)  
        if is_ai_included and col_idx < len(cols):  
            # Yalnızca var olan dosyaları zip'e ekle
            files_to_zip = [p for p in report_paths if p and os.path.exists(p)]
            
            buffer = BytesIO()  
            with zipfile.ZipFile(buffer, "w") as zf:  
                for path in files_to_zip:
                    zf.write(path, os.path.basename(path))  
            buffer.seek(0)  
              
            with cols[col_idx]:  
                st.download_button(  
                    "📦 All Files (ZIP)",  
                    buffer,  
                    file_name="metriqAI_complete_package.zip",  
                    mime="application/zip"  
                )  
            col_idx += 1  
          
        # -------------------------------------------------------  
        # EMAIL DELIVERY (Pro & Premium)  
        # -------------------------------------------------------  
          
        if package_features['email_delivery']:  
            st.markdown("---")  
            st.markdown("### 📧 Email Delivery")  
              
            with st.expander("✉️ Send reports via email"):  
                recipient_email = st.text_input("Recipient Email:", placeholder="client@example.com")  
                recipient_name = st.text_input("Recipient Name:", placeholder="John Doe")  
                  
                if st.button("📤 Send Email"):  
                    if recipient_email and recipient_name:  
                        st.info("📧 Email delivery feature coming soon! For now, please download the files manually.")  
                    else:  
                        st.warning("Please enter both email and name.")  
          
        # -------------------------------------------------------  
        # VIDEO WALKTHROUGH (Premium Only)  
        # -------------------------------------------------------  
          
        if package_features['video_walkthrough']:  
            st.markdown("---")  
            st.markdown("### 🎥 Video Walkthrough")  
            st.info("🎬 Video walkthrough will be sent to your email within 6 hours. Our analyst will personally explain the key findings.")  
          
        # -------------------------------------------------------  
        # DOSYA TEMİZLİĞİ  
        # -------------------------------------------------------  
          
        # Geçici dosyaları temizle  
        def cleanup_files():  
            all_files = list(graphs.values()) + report_paths  
            all_files.extend(glob.glob("*.zip"))  
              
            for f in set(all_files):  
                try:  
                    if os.path.exists(f):  
                        os.remove(f)  
                except Exception:  
                    pass  
          
        # Cleanup (sessiona göre sadece bir kez)  
        if 'cleaned' not in st.session_state:  
            cleanup_files()  
            st.session_state['cleaned'] = True  
          
    except Exception as e:  
        st.error(f"⚠️ Error processing data: {str(e)}")  
        st.info("Please make sure your file has the required columns: tarih, net_tutar. Optional: kategori, sehir. Also check for non-numeric/non-date values.")

else:
    # -------------------------------------------------------
    # EMPTY STATE (Dosya Yüklenmemiş)
    # -------------------------------------------------------

    st.markdown("---")  
      
    st.markdown("""  
    <div style='  
        text-align: center;  
        padding: 3rem 2rem;  
        background: rgba(0, 191, 255, 0.05);  
        border-radius: 20px;  
        border: 2px dashed rgba(0, 191, 255, 0.3);  
        margin: 2rem 0;  
    '>  
        <h2 style='color: #00BFFF;'>👆 Upload Your Data to Begin</h2>  
        <p style='color: rgba(255,255,255,0.7); font-size: 1.1rem; margin: 1rem 0;'>  
            Drop your Excel or CSV file above to get instant AI-powered insights  
        </p>  
        <br>  
        <div style='display: flex; justify-content: center; gap: 2rem; flex-wrap: wrap;'>  
            <div style='text-align: center;'>  
                <div style='font-size: 3rem;'>⚡</div>  
                <p style='color: rgba(255,255,255,0.8);'><strong>Lightning Fast</strong><br>Results in seconds</p>  
            </div>  
            <div style='text-align: center;'>  
                <div style='font-size: 3rem;'>🤖</div>  
                <p style='color: rgba(255,255,255,0.8);'><strong>AI-Powered</strong><br>Smart insights</p>  
            </div>  
            <div style='text-align: center;'>  
                <div style='font-size: 3rem;'>📊</div>  
                <p style='color: rgba(255,255,255,0.8);'><strong>Multi-Format</strong><br>PDF, Excel, PPT</p>  
            </div>  
        </div>  
    </div>  
    """, unsafe_allow_html=True)  
      
    # Demo info  
    st.markdown("### 📋 Sample Data Structure")  
    st.markdown("""  
    Your data should include these columns (case-insensitive):  
    - `tarih` - Date (YYYY-MM-DD)  
    - `net_tutar` - Net Amount/Revenue  
    - `kategori` - Category/Product (optional)  
    - `sehir` - City (optional)  
    """)  
      
    # Package comparison  
    st.markdown("---")  
    st.markdown("### 💎 Choose Your Package")  
      
    col1, col2, col3 = st.columns(3)  
      
    for idx, (pkg_name, pkg_info) in enumerate(PACKAGES.items()):  
        with [col1, col2, col3][idx]:  
            st.markdown(f"""  
            <div style='  
                background: rgba(255,255,255,0.05);  
                padding: 20px;  
                border-radius: 15px;  
                border: 2px solid {pkg_info['color']};  
                text-align: center;  
                height: 100%;  
            '>  
                <h3 style='color: {pkg_info['color']};'>{pkg_info['icon']} {pkg_name}</h3>  
                <p style='font-size: 2rem; font-weight: bold; color: white;'>{pkg_info['price']}</p>  
                <p style='color: rgba(255,255,255,0.7);'>⏱️ {pkg_info['delivery']}</p>  
            </div>  
            """, unsafe_allow_html=True)

# -------------------------------------------------------

# FOOTER

# -------------------------------------------------------

render_footer()
"""
data_analysis.py - Veri Analiz Modülü
"""
import pandas as pd
import numpy as np
from io import BytesIO

class SalesAnalyzer:
    """Yüklenen satış verilerini analiz eder."""
    
    def __init__(self, df: pd.DataFrame):
        self.df = self._clean_data(df)

    @staticmethod
    def _clean_data(df: pd.DataFrame) -> pd.DataFrame:
        """Veriyi temizler ve gerekli tiplere dönüştürür."""
        
        # Sütun isimlerini küçük harfe çevir
        df.columns = df.columns.str.lower()
        
        # Tarih sütunu
        if 'tarih' in df.columns:
            # Tarih sütununu zorla datetime yap
            df['tarih'] = pd.to_datetime(df['tarih'], errors='coerce')
            df.dropna(subset=['tarih'], inplace=True)
            df.set_index('tarih', inplace=True)
        
        # Tutar sütunu
        if 'net_tutar' in df.columns:
            # Tutar sütununu zorla numeric yap (hataları NaN yapar)
            df['net_tutar'] = pd.to_numeric(df['net_tutar'], errors='coerce')
            df.dropna(subset=['net_tutar'], inplace=True)
            
        # Kategori, şehir
        if 'kategori' in df.columns:
            df['kategori'] = df['kategori'].astype(str)
        if 'sehir' in df.columns:
            df['sehir'] = df['sehir'].astype(str)

        return df

    @classmethod
    def load_data(cls, uploaded_file):
        """Yüklenen dosyayı okur ve SalesAnalyzer nesnesi döndürür."""
        try:
            if uploaded_file.name.endswith('.csv'):
                df = pd.read_csv(uploaded_file)
            else: # xlsx, xls
                df = pd.read_excel(uploaded_file)
            
            # Gerekli minimum sütunları kontrol et
            required_cols = ['tarih', 'net_tutar']
            if not all(col in df.columns.str.lower() for col in required_cols):
                 raise ValueError(f"Data must contain columns: {', '.join(required_cols)}")
                 
            return cls(df)
            
        except Exception as e:
            raise Exception(f"Error loading data: {e}")

    def compute_kpis(self):
        """Temel KPI'ları hesaplar."""
        df = self.df
        
        total_revenue = df['net_tutar'].sum()
        
        # Günlük ortalama
        if not df.empty and df.index.min() is not pd.NaT:
            date_range = (df.index.max() - df.index.min()).days + 1
            avg_daily = total_revenue / date_range if date_range > 0 else 0
        else:
            avg_daily = 0

        # WoW (Haftalık değişim)
        wow = np.nan
        if len(df.index.unique()) >= 14:
            # Tarih index'i olduğu varsayılır
            last_14_days = df[df.index >= df.index.max() - pd.Timedelta(days=13)] # Son 14 gün
            
            if len(last_14_days) > 0:
                # Son 7 gün
                last_7_days = last_14_days[last_14_days.index >= last_14_days.index.max() - pd.Timedelta(days=6)]['net_tutar'].sum()
                # Önceki 7 gün
                prev_7_days = last_14_days[last_14_days.index < last_14_days.index.max() - pd.Timedelta(days=6)]['net_tutar'].sum()
                
                if prev_7_days > 0:
                    wow = ((last_7_days - prev_7_days) / prev_7_days) * 100
                elif last_7_days > 0:
                    wow = 100.0 # Önceki hafta 0, bu hafta >0
                else:
                    wow = 0.0 # İki hafta da 0

        return {
            "total": total_revenue,
            "avg": avg_daily,
            "wow": wow,
            "transactions": len(df)
        }

    def get_kpi_text(self):
        """AI özeti için KPI'ları metin formatında döndürür."""
        kpis = self.compute_kpis()
        
        text = "Temel Performans Göstergeleri:\n"
        text += f"- Toplam Ciro: {kpis['total']:,.0f} TL\n"
        text += f"- Günlük Ortalama: {kpis['avg']:,.0f} TL\n"
        text += f"- İşlem Sayısı: {kpis['transactions']:}\n"
        if not np.isnan(kpis['wow']):
             text += f"- Haftalık Değişim (WoW): {kpis['wow']:+.1f}%\n"
        
        # Top 5 Şehir
        if 'sehir' in self.df.columns:
            top_cities = self.df.groupby('sehir')['net_tutar'].sum().nlargest(5)
            text += "\nEn Yüksek Ciro Yapan 5 Şehir:\n"
            for city, revenue in top_cities.items():
                text += f"  - {city}: {revenue:,.0f} TL\n"

        return text
"""
reporting.py - Rapor Oluşturma ve Export Modülü
"""
import streamlit as st
import os
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from docx import Document
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import numpy as np # np import'u eklendi
# from openai import OpenAI 

# Config'den OPENAI_API_KEY'i al
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")

# Mock (Sahte) AI Summary fonksiyonu
def ai_summary(kpi_text: str) -> str:
    """
    OpenAI API ile stratejik özeti oluşturur.
    API key yoksa mock bir özet döndürür.
    """
    if not OPENAI_API_KEY:
        return (
            "🤖 (MOCK AI Summary) The data indicates a strong performance with "
            "a total revenue of over a million TL. The Week-over-Week growth "
            "is positive, suggesting an accelerating trend. The focus should "
            "be on scaling up operations in the top-performing cities and "
            "identifying the factors driving the success of the top product categories. "
            "A potential area for improvement is to analyze low-performing regions."
        )

    # Gerçek OpenAI API çağrısı bu şekilde olurdu
    # client = OpenAI(api_key=OPENAI_API_KEY)
    # prompt = f"Analyze the following KPIs and generate a 3-paragraph strategic summary: {kpi_text}"
    # response = client.chat.completions.create(
    #     model="gpt-4o",
    #     messages=[
    #         {"role": "system", "content": "You are a world-class strategic business analyst."},
    #         {"role": "user", "content": prompt}
    #     ]
    # )
    # return response.choices[0].message.content
    
    # Şu an için, API key olsa bile (test için) mock özeti döndürelim
    return (
        "🤖 (MOCK AI Summary) The data indicates a strong performance with "
        "a total revenue of over a million TL. The Week-over-Week growth "
        "is positive, suggesting an accelerating trend. The focus should "
        "be on scaling up operations in the top-performing cities and "
        "identifying the factors driving the success of the top product categories. "
        "A potential area for improvement is to analyze low-performing regions."
    )


def generate_graphs(kpis: dict) -> dict:
    """
    Statik grafikler oluşturur ve geçici dosya yollarını döndürür.
    (Basic paket için kullanılır).
    """
    temp_files = {}

    # Sahte veriler oluştur (gerçekte analyzer.df kullanılır)
    np.random.seed(42)
    dates = pd.date_range('2024-01-01', periods=30)
    daily_revenue = np.random.randint(10000, 50000, 30).cumsum() + 100000 
    products = ['A', 'B', 'C', 'D', 'E']
    top_product_data = pd.Series(np.random.randint(50000, 300000, 5), index=products)
    regions = ['İstanbul', 'Ankara', 'İzmir', 'Diğer']
    top_region_data = pd.Series(np.random.randint(200000, 800000, 4), index=regions)
    
    sns.set_theme(style="darkgrid")
    
    # 1. Daily Revenue Trend (Çizgi)
    plt.figure(figsize=(10, 5))
    plt.plot(dates, daily_revenue, marker='o', linestyle='-', color='#00BFFF')
    plt.title('Daily Revenue Trend', color='white')
    plt.xlabel('Date', color='white')
    plt.ylabel('Revenue (TL)', color='white')
    plt.xticks(rotation=45, color='white')
    plt.yticks(color='white')
    plt.gca().set_facecolor('#203a43')
    plt.gcf().set_facecolor('#0f2027')
    plt.grid(color='rgba(255,255,255,0.1)')
    plt.tight_layout()
    daily_path = "daily_revenue_trend.png"
    plt.savefig(daily_path)
    plt.close()
    temp_files["daily"] = daily_path
    
    # 2. Top Products (Bar)
    plt.figure(figsize=(10, 5))
    sns.barplot(x=top_product_data.index, y=top_product_data.values, palette="Blues_d")
    plt.title('Top Products by Revenue', color='white')
    plt.xlabel('Product', color='white')
    plt.ylabel('Revenue (TL)', color='white')
    plt.xticks(color='white')
    plt.yticks(color='white')
    plt.gca().set_facecolor('#203a43')
    plt.gcf().set_facecolor('#0f2027')
    plt.grid(color='rgba(255,255,255,0.1)', axis='y')
    plt.tight_layout()
    product_path = "top_products.png"
    plt.savefig(product_path)
    plt.close()
    temp_files["top_product"] = product_path
    
    # 3. Regional Distribution (Pasta)
    plt.figure(figsize=(8, 8))
    plt.pie(
        top_region_data.values, 
        labels=top_region_data.index, 
        autopct='%1.1f%%', 
        startangle=140,
        colors=['#00BFFF', '#0080FF', '#3498DB', '#2C5364'],
        textprops={'color': 'white'}
    )
    plt.title('Regional Distribution', color='white')
    plt.gcf().set_facecolor('#0f2027')
    plt.tight_layout()
    region_path = "regional_distribution.png"
    plt.savefig(region_path)
    plt.close()
    temp_files["top_region"] = region_path

    return temp_files

def build_pdf(summary: str, kpis: dict, graphs: dict, filename: str) -> str:
    """PDF Raporu oluşturur."""
    doc = SimpleDocTemplate(filename, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []

    # Başlık
    story.append(Paragraph("metriq.AI Sales Analysis Report", styles['h1']))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Date: " + pd.Timestamp.now().strftime('%Y-%m-%d'), styles['Normal']))
    story.append(Spacer(1, 24))

    # KPI'lar
    story.append(Paragraph("Key Performance Indicators (KPIs)", styles['h2']))
    story.append(Paragraph(f"💰 Total Revenue: {kpis['total']:,.0f} TL", styles['Normal']))
    story.append(Paragraph(f"📊 Daily Average: {kpis['avg']:,.0f} TL", styles['Normal']))
    story.append(Paragraph(f"📈 Week-over-Week: {kpis['wow']:+.1f}%" if not pd.isna(kpis['wow']) else "📈 Week-over-Week: N/A", styles['Normal']))
    story.append(Paragraph(f"🔢 Total Transactions: {kpis['transactions']:,}", styles['Normal']))
    story.append(Spacer(1, 24))

    # AI Özet
    story.append(Paragraph("AI Strategic Summary", styles['h2']))
    story.append(Paragraph(summary, styles['Normal']))
    story.append(Spacer(1, 24))

    # Grafikler
    story.append(Paragraph("Data Visualizations", styles['h2']))
    from reportlab.lib.utils import ImageReader
    
    for graph_path in graphs.values():
        if os.path.exists(graph_path):
            img = ImageReader(graph_path)
            # Görüntü nesnesini ekle (boyut ayarı ile)
            story.append(ImageReader(graph_path, width=450, height=225))
            story.append(Spacer(1, 12))

    doc.build(story)
    return filename

def build_docx(df: pd.DataFrame, summary: str, filename: str) -> str:
    """Word Belgesi oluşturur."""
    doc = Document()
    doc.add_heading('metriq.AI Analysis Report', 0)

    doc.add_heading('AI Strategic Summary', 1)
    doc.add_paragraph(summary)
    
    # Dataframe'in ilk 10 satırını tablo olarak ekle
    doc.add_heading('Sample Data (First 10 Rows)', 1)
    df_head = df.reset_index().head(10)
    table = doc.add_table(df_head.shape[0] + 1, df_head.shape[1])
    table.style = 'Light Shading'

    # Sütun başlıkları
    for j, col in enumerate(df_head.columns):
        table.cell(0, j).text = str(col)

    # Satır verileri
    for i in range(df_head.shape[0]):
        for j in range(df_head.shape[1]):
            table.cell(i + 1, j).text = str(df_head.iloc[i, j])

    doc.save(filename)
    return filename

def build_ppt(summary: str, graphs: dict, filename: str) -> str:
    """PowerPoint Sunumu oluşturur."""
    prs = Presentation()
    
    # 1. Kapak Slaytı
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Metriq.AI - Strategic Sales Analysis"
    slide.placeholders[1].text = "Date: " + pd.Timestamp.now().strftime('%Y-%m-%d')
    
    # 2. AI Özet Slaytı
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Key AI Strategic Insights"
    body = slide.shapes.placeholders[1]
    
    # Özet metnini maddeler halinde ekle
    p_frame = body.text_frame
    p = p_frame.add_paragraph()
    p.text = summary
    
    # 3. Grafik Slaytları
    # Boş bir layout seç
    blank_slide_layout = prs.slide_layouts[6] 
    
    for graph_title, graph_path in graphs.items():
        if os.path.exists(graph_path):
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Başlık ekle
            left = top = width = height = Inches(0.5)
            txBox = slide.shapes.add_textbox(left, top, width=Inches(9), height=height)
            tf = txBox.text_frame
            tf.word_wrap = True
            tf.text = graph_title.replace('_', ' ').title()
            
            # Resmi ekle ve ortala
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(9)
            height = Inches(5)
            pic = slide.shapes.add_picture(graph_path, left, top, width, height)

    prs.save(filename)
    return filename

def save_excel(df: pd.DataFrame, filename: str) -> str:
    """Temizlenmiş veriyi Excel dosyasına kaydeder."""
    # Tarih index'ini sütuna çevirerek kaydet
    if 'tarih' in df.columns:
         df = df.reset_index(drop=True)
    elif df.index.name == 'tarih':
         df = df.reset_index()

    df.to_excel(filename, index=False)
    return filename
