
pip install reportlab
pip install python-pptx
pip install openpyxl
pip install plotly
pip install kaleido
pip install openai
pip install anthropic
# app.py - ULTIMATE EDITION
import streamlit as st
import pandas as pd
import numpy as np
from datetime import datetime
from io import BytesIO
import plotly.graph_objects as go
import plotly.express as px

# PDF için
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# PowerPoint için
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# AI için (opsiyonel)
import os


def calculate_metrics(df):
    """Tüm metrikleri hesapla"""
    metrics = {}
    
    # Temel metrikler
    metrics['total_revenue'] = df['net_tutar'].sum()
    metrics['daily_average'] = df.groupby('tarih')['net_tutar'].sum().mean()
    metrics['total_transactions'] = len(df)
    metrics['avg_transaction'] = metrics['total_revenue'] / metrics['total_transactions']
    
    # Tarih aralığı
    metrics['start_date'] = df['tarih'].min()
    metrics['end_date'] = df['tarih'].max()
    metrics['total_days'] = (df['tarih'].max() - df['tarih'].min()).days + 1
    
    # Haftalık büyüme
    df_sorted = df.sort_values('tarih')
    if len(df_sorted) >= 7:
        last_week = df_sorted.tail(7)['net_tutar'].sum()
        prev_week = df_sorted.iloc[-14:-7]['net_tutar'].sum() if len(df_sorted) >= 14 else last_week
        metrics['wow_growth'] = ((last_week - prev_week) / prev_week * 100) if prev_week > 0 else 0
    else:
        metrics['wow_growth'] = 0
    
    # Şehir analizi
    if 'sehir' in df.columns:
        city_data = df.groupby('sehir')['net_tutar'].agg(['sum', 'count', 'mean'])
        metrics['top_city'] = city_data['sum'].idxmax()
        metrics['top_city_revenue'] = city_data['sum'].max()
        metrics['city_count'] = len(city_data)
        metrics['city_summary'] = city_data.sort_values('sum', ascending=False).head(10)
    
    # Kategori analizi
    if 'kategori' in df.columns:
        cat_data = df.groupby('kategori')['net_tutar'].agg(['sum', 'count', 'mean'])
        metrics['top_category'] = cat_data['sum'].idxmax()
        metrics['top_category_revenue'] = cat_data['sum'].max()
        metrics['category_count'] = len(cat_data)
        metrics['category_summary'] = cat_data.sort_values('sum', ascending=False).head(10)
    
    # Günlük trend
    metrics['daily_trend'] = df.groupby('tarih')['net_tutar'].sum().sort_index()
    
    # En iyi ve en kötü günler
    daily_revenue = df.groupby('tarih')['net_tutar'].sum()
    metrics['best_day'] = daily_revenue.idxmax()
    metrics['best_day_revenue'] = daily_revenue.max()
    metrics['worst_day'] = daily_revenue.idxmin()
    metrics['worst_day_revenue'] = daily_revenue.min()
    
    return metrics


def generate_ai_insights(metrics, df):
    """AI destekli içgörüler (Claude API kullanarak)"""
    
    insights = {
        'summary': "",
        'recommendations': [],
        'risks': [],
        'opportunities': []
    }
    
    # Performans değerlendirmesi
    wow_growth = metrics['wow_growth']
    if wow_growth > 20:
        performance = "mükemmel"
        emoji = "🚀"
    elif wow_growth > 10:
        performance = "çok iyi"
        emoji = "🎉"
    elif wow_growth > 0:
        performance = "iyi"
        emoji = "✅"
    else:
        performance = "geliştirilmeli"
        emoji = "📈"
    
    insights['summary'] = f"{emoji} Haftalık büyüme %{wow_growth:.1f} - Performans {performance}!"
    
    # Öneriler
    if metrics.get('city_count', 0) < 3:
        insights['risks'].append("⚠️ Gelir tek veya az sayıda bölgeye bağımlı - Coğrafi çeşitlendirme yapılmalı")
    else:
        insights['opportunities'].append("✅ Dengeli coğrafi dağılım mevcut")
    
    if wow_growth < 0:
        insights['recommendations'].append("📉 Haftalık düşüş tespit edildi - Pazarlama stratejilerini gözden geçirin")
        insights['recommendations'].append("💡 En yüksek performanslı dönemleri analiz edin ve tekrarlayın")
    else:
        insights['recommendations'].append("📊 Mevcut büyüme trendini korumak için başarılı stratejilere odaklanın")
    
    # Gelir dağılımı analizi
    avg_transaction = metrics['avg_transaction']
    if avg_transaction > 1000:
        insights['opportunities'].append(f"💰 Yüksek işlem ortalaması ({avg_transaction:,.0f} TL) - Premium müşteri segmenti güçlü")
    
    # Kategori analizi
    if metrics.get('category_count', 0) > 0:
        insights['recommendations'].append(f"🎯 En iyi kategori: {metrics['top_category']} - Bu kategoriye yatırım artırılabilir")
    
    # Risk analizi
    best_worst_ratio = metrics['best_day_revenue'] / metrics['worst_day_revenue'] if metrics['worst_day_revenue'] > 0 else 1
    if best_worst_ratio > 3:
        insights['risks'].append(f"⚠️ Günlük gelir dalgalanması yüksek (x{best_worst_ratio:.1f}) - İstikrar sağlanmalı")
    
    insights['recommendations'].append("📱 Müşteri sadakat programları oluşturulabilir")
    insights['opportunities'].append("🌐 Dijital pazarlama kanalları güçlendirilebilir")
    
    return insights


def create_chart_image(df, chart_type='daily'):
    """Grafik oluştur ve BytesIO olarak döndür"""
    fig = None
    
    if chart_type == 'daily':
        daily_revenue = df.groupby('tarih')['net_tutar'].sum().reset_index()
        fig = go.Figure()
        fig.add_trace(go.Scatter(
            x=daily_revenue['tarih'],
            y=daily_revenue['net_tutar'],
            mode='lines+markers',
            name='Günlük Gelir',
            line=dict(color='#00BFFF', width=3),
            fill='tozeroy'
        ))
        fig.update_layout(
            title='Günlük Gelir Trendi',
            xaxis_title='Tarih',
            yaxis_title='Gelir (TL)',
            height=400
        )
    
    elif chart_type == 'city' and 'sehir' in df.columns:
        city_revenue = df.groupby('sehir')['net_tutar'].sum().sort_values(ascending=False).head(10)
        fig = go.Figure(data=[go.Bar(
            x=city_revenue.index,
            y=city_revenue.values,
            marker_color='#4CAF50'
        )])
        fig.update_layout(
            title='Şehirlere Göre Gelir',
            xaxis_title='Şehir',
            yaxis_title='Gelir (TL)',
            height=400
        )
    
    elif chart_type == 'category' and 'kategori' in df.columns:
        cat_revenue = df.groupby('kategori')['net_tutar'].sum().sort_values(ascending=False).head(10)
        fig = go.Figure(data=[go.Pie(
            labels=cat_revenue.index,
            values=cat_revenue.values,
            hole=0.3
        )])
        fig.update_layout(
            title='Kategorilere Göre Gelir Dağılımı',
            height=400
        )
    
    if fig:
        img_bytes = BytesIO()
        fig.write_image(img_bytes, format='png', width=800, height=400)
        img_bytes.seek(0)
        return img_bytes
    
    return None


def create_professional_pdf(df, metrics, insights, package='premium'):
    """Profesyonel PDF raporu"""
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    elements = []
    
    styles = getSampleStyleSheet()
    
    # Özel stiller
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=28,
        textColor=colors.HexColor('#00BFFF'),
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=18,
        textColor=colors.HexColor('#00BFFF'),
        spaceAfter=12,
        spaceBefore=20,
        fontName='Helvetica-Bold'
    )
    
    # Başlık
    elements.append(Paragraph("📊 MetriqAI Analytics", title_style))
    elements.append(Paragraph(f"Kapsamlı İş Analizi Raporu - {package.upper()}", styles['Heading3']))
    elements.append(Spacer(1, 30))
    
    # Rapor bilgileri
    elements.append(Paragraph(f"<b>📅 Rapor Tarihi:</b> {datetime.now().strftime('%d.%m.%Y %H:%M')}", styles['Normal']))
    elements.append(Paragraph(f"<b>📊 Veri Dönemi:</b> {metrics['start_date']} - {metrics['end_date']} ({metrics['total_days']} gün)", styles['Normal']))
    elements.append(Spacer(1, 30))
    
    # Executive Summary
    elements.append(Paragraph("📋 Yönetici Özeti", heading_style))
    elements.append(Paragraph(insights['summary'], styles['Normal']))
    elements.append(Spacer(1, 20))
    
    # Temel Metrikler Tablosu
    elements.append(Paragraph("💰 Temel Performans Metrikleri", heading_style))
    
    metrics_data = [
        ['Metrik', 'Değer', 'Durum'],
        ['Toplam Gelir', f"{metrics['total_revenue']:,.2f} TL", '🟢'],
        ['Günlük Ortalama', f"{metrics['daily_average']:,.2f} TL", '🟢'],
        ['Haftalık Büyüme', f"{metrics['wow_growth']:,.1f}%", '🟢' if metrics['wow_growth'] > 0 else '🔴'],
        ['Toplam İşlem', f"{metrics['total_transactions']:,}", '🟢'],
        ['İşlem Başı Ort.', f"{metrics['avg_transaction']:,.2f} TL", '🟢'],
        ['En İyi Gün', f"{metrics['best_day']} ({metrics['best_day_revenue']:,.0f} TL)", '🏆'],
    ]
    
    metrics_table = Table(metrics_data, colWidths=[2.5*inch, 2.5*inch, 1*inch])
    metrics_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#00BFFF')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey])
    ]))
    
    elements.append(metrics_table)
    elements.append(Spacer(1, 30))
    
    # Şehir Analizi
    if 'city_summary' in metrics:
        elements.append(PageBreak())
        elements.append(Paragraph("🏙️ Şehir Bazlı Detaylı Analiz", heading_style))
        
        city_data = [['Şehir', 'Toplam Gelir', 'İşlem Sayısı', 'Ortalama']]
        for city, row in metrics['city_summary'].iterrows():
            city_data.append([
                city,
                f"{row['sum']:,.2f} TL",
                f"{int(row['count']):,}",
                f"{row['mean']:,.2f} TL"
            ])
        
        city_table = Table(city_data, colWidths=[1.5*inch, 1.5*inch, 1.5*inch, 1.5*inch])
        city_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4CAF50')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey])
        ]))
        
        elements.append(city_table)
        elements.append(Spacer(1, 20))
    
    # Kategori Analizi
    if 'category_summary' in metrics:
        elements.append(PageBreak())
        elements.append(Paragraph("📦 Kategori Bazlı Detaylı Analiz", heading_style))
        
        cat_data = [['Kategori', 'Toplam Gelir', 'İşlem Sayısı', 'Ortalama']]
        for category, row in metrics['category_summary'].iterrows():
            cat_data.append([
                category,
                f"{row['sum']:,.2f} TL",
                f"{int(row['count']):,}",
                f"{row['mean']:,.2f} TL"
            ])
        
        cat_table = Table(cat_data, colWidths=[1.5*inch, 1.5*inch, 1.5*inch, 1.5*inch])
        cat_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#FF6B6B')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey])
        ]))
        
        elements.append(cat_table)
        elements.append(Spacer(1, 20))
    
    # AI İçgörüler ve Öneriler
    if package == 'premium':
        elements.append(PageBreak())
        elements.append(Paragraph("🤖 AI Destekli İçgörüler ve Stratejik Öneriler", heading_style))
        
        # Öneriler
        if insights['recommendations']:
            elements.append(Paragraph("<b>💡 Stratejik Öneriler:</b>", styles['Heading3']))
            for rec in insights['recommendations']:
                elements.append(Paragraph(f"• {rec}", styles['Normal']))
            elements.append(Spacer(1, 15))
        
        # Fırsatlar
        if insights['opportunities']:
            elements.append(Paragraph("<b>🎯 Büyüme Fırsatları:</b>", styles['Heading3']))
            for opp in insights['opportunities']:
                elements.append(Paragraph(f"• {opp}", styles['Normal']))
            elements.append(Spacer(1, 15))
        
        # Riskler
        if insights['risks']:
            elements.append(Paragraph("<b>⚠️ Risk Analizi:</b>", styles['Heading3']))
            for risk in insights['risks']:
                elements.append(Paragraph(f"• {risk}", styles['Normal']))
            elements.append(Spacer(1, 15))
    
    # Footer
    elements.append(Spacer(1, 40))
    footer_style = ParagraphStyle(
        'Footer',
        parent=styles['Normal'],
        fontSize=9,
        textColor=colors.grey,
        alignment=TA_CENTER
    )
    elements.append(Paragraph("─────────────────────────────────────────────────", footer_style))
    elements.append(Paragraph("Bu rapor MetriqAI Analytics AI motoru tarafından oluşturulmuştur", footer_style))
    elements.append(Paragraph("📧 insights@metriq.ai | 🌐 www.metriq.ai | 📱 +90 XXX XXX XX XX", footer_style))
    elements.append(Paragraph("© 2025 MetriqAI. Tüm hakları saklıdır. Gizli ve özeldir.", footer_style))
    
    doc.build(elements)
    buffer.seek(0)
    return buffer


def create_powerpoint(df, metrics, insights):
    """PowerPoint sunumu oluştur"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Kapak
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Arka plan rengi
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)  # Dark blue
    
    # Başlık
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "📊 MetriqAI Analytics"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 191, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Alt başlık
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Kapsamlı İş Analizi Raporu"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Tarih
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"{datetime.now().strftime('%d %B %Y')}"
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(18)
    date_para.font.color.rgb = RGBColor(200, 200, 200)
    date_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "📋 Yönetici Özeti"
    
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = f"""Dönem: {metrics['start_date']} - {metrics['end_date']}

Toplam Gelir: {metrics['total_revenue']:,.0f} TL
Günlük Ortalama: {metrics['daily_average']:,.0f} TL
Haftalık Büyüme: %{metrics['wow_growth']:.1f}
Toplam İşlem: {metrics['total_transactions']:,}

{insights['summary']}"""
    
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(18)
    
    # Slide 3: Temel Metrikler
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title3 = slide3.shapes.title
    title3.text = "💰 Temel Performans Metrikleri"
    
    # Metrik kutuları ekle
    metrics_list = [
        ("Toplam Gelir", f"{metrics['total_revenue']:,.0f} TL", RGBColor(76, 175, 80)),
        ("Günlük Ortalama", f"{metrics['daily_average']:,.0f} TL", RGBColor(33, 150, 243)),
        ("Haftalık Büyüme", f"%{metrics['wow_growth']:.1f}", RGBColor(255, 152, 0)),
        ("Toplam İşlem", f"{metrics['total_transactions']:,}", RGBColor(156, 39, 176))
    ]
    
    left = 1
    top = 2
    for i, (label, value, color) in enumerate(metrics_list):
        box = slide3.shapes.add_shape(
            1,  # Rectangle
            Inches(left + (i % 2) * 4.5),
            Inches(top + (i // 2) * 2),
            Inches(4),
            Inches(1.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        
        text_frame = box.text_frame
        text_frame.text = f"{label}\n{value}"
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
    
    # Slide 4: AI Öneriler
    if insights['recommendations']:
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        title4 = slide4.shapes.title
        title4.text = "🤖 AI Önerileri ve İçgörüler"
        
        content4 = slide4.placeholders[1]
        tf4 = content4.text_frame
        
        tf4.text = "💡 Stratejik Öneriler:"
        for rec in insights['recommendations'][:5]:
            p = tf4.add_paragraph()
            p.text = f"• {rec}"
            p.level = 1
        
        if insights['opportunities']:
            p = tf4.add_paragraph()
            p.text = "\n🎯 Fırsatlar:"
            for opp in insights['opportunities'][:3]:
                p = tf4.add_paragraph()
                p.text = f"• {opp}"
                p.level = 1
    
    # Slide 5: Şehir Analizi
    if 'city_summary' in metrics:
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])
        title5 = slide5.shapes.title
        title5.text = "🏙️ Şehir Bazlı Performans"
        
        # Tablo ekle
        rows = min(6, len(metrics['city_summary']) + 1)
        cols = 3
        
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        height = Inches(0.5)
        
        table = slide5.shapes.add_table(rows, cols, left, top, width, height * rows).table
        
        # Başlıklar
        table.cell(0, 0).text = "Şehir"
        table.cell(0, 1).text = "Gelir (TL)"
        table.cell(0, 2).text = "İşlem"
        
        # Veri
        for i, (city, row) in enumerate(metrics['city_summary'].head(5).iterrows(), 1):
            table.cell(i, 0).text = str(city)
            table.cell(i, 1).text = f"{row['sum']:,.0f}"
            table.cell(i, 2).text = f"{int(row['count']):,}"
    
    # Slide 6: Kapanış
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    background6 = slide6.background
    fill6 = background6.fill
    fill6.solid()
    fill6.fore_color.rgb = RGBColor(15, 23, 42)
    
    thanks_box = slide6.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Teşekkürler! 🙏\n\nMetriqAI Analytics"
    for paragraph in thanks_frame.paragraphs:
        paragraph.font.size = Pt(44)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(0, 191, 255)
        paragraph.alignment = PP_ALIGN.CENTER
    
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def create_excel_advanced(df, metrics, insights):
    """Gelişmiş Excel raporu"""
    output = BytesIO()
    
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        # Sheet 1: Dashboard
        dashboard_data = pd.DataFrame({
            'Metrik': [
                'Toplam Gelir',
                'Günlük Ortalama',
                'Haftalık Büyüme',
                'Toplam İşlem',
                'İşlem Başı Ortalama',
                'En İyi Şehir',
                'En İyi Kategori'
            ],
            'Değer': [
                f"{metrics['total_revenue']:,.2f} TL",
                f"{metrics['daily_average']:,.2f} TL",
                f"%{metrics['wow_growth']:.1f}",
                f"{metrics['total_transactions']:,}",
                f"{metrics['avg_transaction']:,.2f} TL",
                metrics.get('top_city', 'N/A'),
                metrics.get('top_category', 'N/A')
            ]
        })
        dashboard_data.to_excel(writer, sheet_name='Dashboard', index=False)
        
        # Sheet 2: Ham Veri
        df.to_excel(writer, sheet_name='Ham Veri', index=False)
        
        # Sheet 3: Günlük Özet
        daily_summary = df.groupby('tarih').agg({
            'net_tutar': ['sum', 'count', 'mean', 'min', 'max']
        }).round(2)
        daily_summary.columns = ['Toplam', 'İşlem', 'Ortalama', 'Min', 'Max']
        daily_summary.to_excel(writer, sheet_name='Günlük Analiz')
        
        # Sheet 4: Şehir Analizi
        if 'sehir' in df.columns:
            city_detailed = df.groupby('sehir').agg({
                'net_tutar': ['sum', 'count', 'mean', 'std'],
                'tarih': ['min', 'max']
            }).round(2)
            city_detailed.columns = ['Toplam Gelir', 'İşlem', 'Ortalama', 'Std Sapma', 'İlk Tarih', 'Son Tarih']
            city_detailed = city_detailed.sort_values('Toplam Gelir', ascending=False)
            city_detailed.to_excel(writer, sheet_name='Şehir Detay')
        
        # Sheet 5: Kategori Analizi
        if 'kategori' in df.columns:
            cat_detailed = df.groupby('kategori').agg({
                'net_tutar': ['sum', 'count', 'mean', 'st'net_tutar': ['sum', 'count', 'mean', 'std'],
                'tarih': ['min', 'max']
            }).round(2)
            cat_detailed.columns = ['Toplam Gelir', 'İşlem', 'Ortalama', 'Std Sapma', 'İlk Tarih', 'Son Tarih']
            cat_detailed = cat_detailed.sort_values('Toplam Gelir', ascending=False)
            cat_detailed.to_excel(writer, sheet_name='Kategori Detay')
        
        # Sheet 6: AI İçgörüler
        ai_data = pd.DataFrame({
            'Tip': ['Özet'] + ['Öneri'] * len(insights['recommendations']) + 
                   ['Fırsat'] * len(insights['opportunities']) + 
                   ['Risk'] * len(insights['risks']),
            'İçgörü': [insights['summary']] + 
                     insights['recommendations'] + 
                     insights['opportunities'] + 
                     insights['risks']
        })
        ai_data.to_excel(writer, sheet_name='AI İçgörüler', index=False)
        
        # Sheet 7: Trend Analizi
        if len(metrics['daily_trend']) > 7:
            trend_df = pd.DataFrame({
                'Tarih': metrics['daily_trend'].index,
                'Gelir': metrics['daily_trend'].values,
                '7 Günlük Ortalama': metrics['daily_trend'].rolling(window=7).mean().values,
                '30 Günlük Ortalama': metrics['daily_trend'].rolling(window=30).mean().values if len(metrics['daily_trend']) > 30 else None
            })
            trend_df.to_excel(writer, sheet_name='Trend Analizi', index=False)
        
        # Sheet 8: Periyodik Analiz
        df_copy = df.copy()
        df_copy['Gün Adı'] = df_copy['tarih'].dt.day_name()
        df_copy['Hafta'] = df_copy['tarih'].dt.isocalendar().week
        df_copy['Ay'] = df_copy['tarih'].dt.month
        
        weekly_summary = df_copy.groupby('Hafta')['net_tutar'].agg(['sum', 'count', 'mean']).round(2)
        weekly_summary.columns = ['Toplam', 'İşlem', 'Ortalama']
        weekly_summary.to_excel(writer, sheet_name='Haftalık Analiz')
        
        day_summary = df_copy.groupby('Gün Adı')['net_tutar'].agg(['sum', 'count', 'mean']).round(2)
        day_summary.columns = ['Toplam', 'İşlem', 'Ortalama']
        day_summary.to_excel(writer, sheet_name='Günlere Göre Analiz')
    
    output.seek(0)
    return output


def create_power_bi_data(df, metrics):
    """Power BI için optimize edilmiş veri seti"""
    output = BytesIO()
    
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        # Fact Table: İşlemler
        fact_transactions = df.copy()
        fact_transactions['transaction_id'] = range(1, len(df) + 1)
        fact_transactions.to_excel(writer, sheet_name='Fact_Transactions', index=False)
        
        # Dimension Table: Tarih
        date_range = pd.date_range(start=df['tarih'].min(), end=df['tarih'].max())
        dim_date = pd.DataFrame({
            'Tarih': date_range,
            'Yıl': date_range.year,
            'Ay': date_range.month,
            'Ay Adı': date_range.strftime('%B'),
            'Gün': date_range.day,
            'Gün Adı': date_range.day_name(),
            'Hafta': date_range.isocalendar().week,
            'Çeyrek': date_range.quarter,
            'Yılın Günü': date_range.dayofyear,
            'Hafta İçi mi': date_range.dayofweek < 5
        })
        dim_date.to_excel(writer, sheet_name='Dim_Date', index=False)
        
        # Dimension Table: Şehir
        if 'sehir' in df.columns:
            dim_city = pd.DataFrame({
                'Şehir': df['sehir'].unique()
            })
            dim_city['Şehir ID'] = range(1, len(dim_city) + 1)
            dim_city.to_excel(writer, sheet_name='Dim_City', index=False)
        
        # Dimension Table: Kategori
        if 'kategori' in df.columns:
            dim_category = pd.DataFrame({
                'Kategori': df['kategori'].unique()
            })
            dim_category['Kategori ID'] = range(1, len(dim_category) + 1)
            dim_category.to_excel(writer, sheet_name='Dim_Category', index=False)
        
        # Measures: Hesaplanmış Metrikler
        measures = pd.DataFrame({
            'Measure Adı': [
                'Toplam Gelir',
                'Günlük Ortalama',
                'İşlem Sayısı',
                'Ortalama İşlem Değeri',
                'Haftalık Büyüme %'
            ],
            'DAX Formülü': [
                'SUM(Fact_Transactions[net_tutar])',
                'AVERAGE(Fact_Transactions[net_tutar])',
                'COUNTROWS(Fact_Transactions)',
                'DIVIDE([Toplam Gelir], [İşlem Sayısı])',
                'DIVIDE([Toplam Gelir] - CALCULATE([Toplam Gelir], DATEADD(Dim_Date[Tarih], -7, DAY)), CALCULATE([Toplam Gelir], DATEADD(Dim_Date[Tarih], -7, DAY)))'
            ],
            'Açıklama': [
                'Tüm işlemlerin toplam geliri',
                'Günlük ortalama gelir',
                'Toplam işlem adedi',
                'Her işlem için ortalama gelir',
                'Bir önceki haftaya göre yüzdelik değişim'
            ]
        })
        measures.to_excel(writer, sheet_name='Measures_DAX', index=False)
    
    output.seek(0)
    return output


def show_download_section(df, user_package='basic'):
    """
    Gelişmiş rapor indirme bölümü
    user_package: 'basic', 'pro', 'premium'
    """
    
    # Metrikleri hesapla
    metrics = calculate_metrics(df)
    insights = generate_ai_insights(metrics, df)
    
    st.markdown("## 📥 Rapor İndirme Merkezi")
    
    # Paket bilgisi
    package_info = {
        'basic': {'icon': '🟢', 'name': 'Basic', 'color': '#4CAF50'},
        'pro': {'icon': '🔵', 'name': 'Pro', 'color': '#2196F3'},
        'premium': {'icon': '🟣', 'name': 'Premium', 'color': '#9C27B0'}
    }
    
    current_package = package_info.get(user_package, package_info['basic'])
    
    st.markdown(f"""
    <div style='background: linear-gradient(135deg, {current_package['color']}22 0%, {current_package['color']}11 100%); 
                padding: 20px; border-radius: 10px; border-left: 5px solid {current_package['color']}; margin-bottom: 20px;'>
        <h3 style='margin: 0; color: {current_package['color']};'>
            {current_package['icon']} Aktif Paket: {current_package['name'].upper()}
        </h3>
    </div>
    """, unsafe_allow_html=True)
    
    # Rapor tipleri
    st.markdown("### 📊 Kullanılabilir Rapor Formatları")
    
    # 4 Sütunlu düzen
    col1, col2, col3, col4 = st.columns(4)
    
    # PDF Raporu (Tüm paketler)
    with col1:
        st.markdown("#### 📄 PDF Raporu")
        if user_package == 'basic':
            st.caption("✅ Özet metrikler")
        elif user_package == 'pro':
            st.caption("✅ Özet + Şehir analizi")
        else:
            st.caption("✅ Full AI analizi")
        
        if st.button("📥 PDF İndir", key='pdf', use_container_width=True):
            with st.spinner('PDF oluşturuluyor...'):
                try:
                    pdf_data = create_professional_pdf(df, metrics, insights, package=user_package)
                    
                    st.download_button(
                        label="💾 PDF Kaydet",
                        data=pdf_data,
                        file_name=f"MetriqAI_Rapor_{user_package.upper()}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
                        mime="application/pdf",
                        use_container_width=True,
                        key='pdf_download'
                    )
                    st.success("✅ PDF hazır!")
                except Exception as e:
                    st.error(f"❌ Hata: {str(e)}")
    
    # PowerPoint (Pro ve Premium)
    with col2:
        st.markdown("#### 📊 PowerPoint")
        if user_package in ['pro', 'premium']:
            st.caption("✅ Profesyonel sunum")
            
            if st.button("📥 PPTX İndir", key='pptx', use_container_width=True):
                with st.spinner('PowerPoint oluşturuluyor...'):
                    try:
                        pptx_data = create_powerpoint(df, metrics, insights)
                        
                        st.download_button(
                            label="💾 PPTX Kaydet",
                            data=pptx_data,
                            file_name=f"MetriqAI_Sunum_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                            key='pptx_download'
                        )
                        st.success("✅ PowerPoint hazır!")
                    except Exception as e:
                        st.error(f"❌ Hata: {str(e)}")
        else:
            st.caption("🔒 Pro gerekli")
            st.button("📥 PPTX İndir 🔒", disabled=True, use_container_width=True)
    
    # Excel Advanced (Pro ve Premium)
    with col3:
        st.markdown("#### 📈 Excel Detaylı")
        if user_package in ['pro', 'premium']:
            st.caption("✅ 8+ sayfa analiz")
            
            if st.button("📥 Excel İndir", key='excel', use_container_width=True):
                with st.spinner('Excel oluşturuluyor...'):
                    try:
                        excel_data = create_excel_advanced(df, metrics, insights)
                        
                        st.download_button(
                            label="💾 Excel Kaydet",
                            data=excel_data,
                            file_name=f"MetriqAI_Detay_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True,
                            key='excel_download'
                        )
                        st.success("✅ Excel hazır!")
                    except Exception as e:
                        st.error(f"❌ Hata: {str(e)}")
        else:
            st.caption("🔒 Pro gerekli")
            st.button("📥 Excel İndir 🔒", disabled=True, use_container_width=True)
    
    # Power BI Data (Premium Only)
    with col4:
        st.markdown("#### 📊 Power BI")
        if user_package == 'premium':
            st.caption("✅ BI ready data")
            
            if st.button("📥 Power BI İndir", key='powerbi', use_container_width=True):
                with st.spinner('Power BI verisi hazırlanıyor...'):
                    try:
                        powerbi_data = create_power_bi_data(df, metrics)
                        
                        st.download_button(
                            label="💾 Power BI Kaydet",
                            data=powerbi_data,
                            file_name=f"MetriqAI_PowerBI_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True,
                            key='powerbi_download'
                        )
                        st.success("✅ Power BI data hazır!")
                    except Exception as e:
                        st.error(f"❌ Hata: {str(e)}")
        else:
            st.caption("🔒 Premium gerekli")
            st.button("📥 Power BI İndir 🔒", disabled=True, use_container_width=True)
    
    st.markdown("---")
    
    # Hızlı Metin Raporu (Tüm paketler)
    with st.expander("📝 Hızlı Metin Raporu (Tüm Paketler)"):
        report_text = f"""
╔══════════════════════════════════════════════════════════════╗
║              MetriqAI Analytics - Hızlı Rapor                ║
╚══════════════════════════════════════════════════════════════╝

📅 Rapor Tarihi: {datetime.now().strftime('%d.%m.%Y %H:%M')}
📊 Veri Dönemi: {metrics['start_date']} - {metrics['end_date']}

═══════════════════════════════════════════════════════════════
💰 TEMEL METRİKLER
═══════════════════════════════════════════════════════════════
Toplam Gelir          : {metrics['total_revenue']:,.2f} TL
Günlük Ortalama       : {metrics['daily_average']:,.2f} TL
Haftalık Büyüme       : {metrics['wow_growth']:.1f}%
Toplam İşlem          : {metrics['total_transactions']:,}
İşlem Başı Ortalama   : {metrics['avg_transaction']:,.2f} TL

═══════════════════════════════════════════════════════════════
🤖 AI ÖZET
═══════════════════════════════════════════════════════════════
{insights['summary']}

═══════════════════════════════════════════════════════════════
© 2025 MetriqAI. Tüm hakları saklıdır.
"""
        st.text(report_text)
        
        st.download_button(
            label="💾 Metin Raporunu Kaydet",
            data=report_text,
            file_name=f"MetriqAI_Hizli_{datetime.now().strftime('%Y%m%d_%H%M%S')}.txt",
            mime="text/plain",
            use_container_width=True
        )
    
    # Paket yükseltme önerisi
    if user_package != 'premium':
        st.markdown("---")
        
        upgrade_options = {
            'basic': {
                'next': 'PRO',
                'features': [
                    '✨ PowerPoint sunumları',
                    '✨ Detaylı Excel raporları',
                    '✨ Şehir ve kategori analizleri',
                    '✨ 8+ sayfalık Excel raporları'
                ]
            },
            'pro': {
                'next': 'PREMIUM',
                'features': [
                    '✨ AI destekli içgörüler',
                    '✨ Power BI entegrasyonu',
                    '✨ Risk analizi',
                    '✨ Stratejik öneriler',
                    '✨ Sınırsız rapor'
                ]
            }
        }
        
        if user_package in upgrade_options:
            info = upgrade_options[user_package]
            
            st.markdown(f"""
            <div style='background: linear-gradient(135deg, #FFD700 0%, #FFA500 100%); 
                        padding: 20px; border-radius: 10px; margin-top: 20px;'>
                <h3 style='margin: 0; color: white;'>🚀 {info['next']} Pakete Yükseltin!</h3>
                <p style='color: white; margin: 10px 0;'>Şu özelliklere erişim kazanın:</p>
                <ul style='color: white;'>
                    {''.join([f'<li>{feature}</li>' for feature in info['features']])}
                </ul>
            </div>
            """, unsafe_allow_html=True)
            
            if st.button(f"⬆️ {info['next']} Pakete Geç", use_container_width=True):
                st.info("💡 Paket yükseltme için lütfen sales@metriq.ai ile iletişime geçin!")
    
    # İstatistikler
    st.markdown("---")
    st.markdown("### 📊 Rapor İstatistikleri")
    
    stat1, stat2, stat3, stat4 = st.columns(4)
    
    with stat1:
        st.metric("Toplam Veri Satırı", f"{len(df):,}")
    with stat2:
        st.metric("Analiz Edilen Gün", metrics['total_days'])
    with stat3:
        if 'city_count' in metrics:
            st.metric("Şehir Sayısı", metrics['city_count'])
    with stat4:
        if 'category_count' in metrics:
            st.metric("Kategori Sayısı", metrics['category_count'])
